# Builds "Parking Slides - 23072026.pptx" from the 22 Jul deck by applying the
# Corridor 1+2-only figure set. Targeted python-pptx edits on a copy — the deck carries
# hand edits that a regenerate would destroy, so the generator is deliberately not used.
#
# Scope change behind every number here: Corridor 03 is withheld pending its conceptual
# design, and the Nalbandyan016 segment (14 Zone A spaces, tagged impact=corridor but
# lying wholly outside every corridor boundary) is excluded. All figures are derived
# from the same GeoJSON the web app serves.
import json
import shutil
from pptx import Presentation
from pptx.util import Inches

# Surveyed-area displacement is read from the GeoJSON rather than typed, so the slide
# cannot drift from the metrics run that produced it. These are a different scope from
# the corridor-wide figures below: six sampled areas, measured, not the whole corridor.
with open('app/static/data/wgs84/field-surveys.geojson', encoding='utf-8') as fh:
    DISP = json.load(fh)['displacement']
D_REMOVED = DISP['removed_supply']      # 1,886
D_DISPLACED = DISP['removed_demand']    # 1,120
D_ONSTREET = DISP['absorb_onstreet']    # 391 — unchanged by the design revision
D_RATIO = round(100.0 * D_DISPLACED / D_REMOVED)        # 59%
D_ABSORB = round(100.0 * D_ONSTREET / D_DISPLACED)      # 35%
YARD_DEP = 92   # residential-yard share of off-street stock; independent of kerb removal

# Per-area rows for slide 4, read from the same areaStats the app's dashboard uses.
with open('app/static/data/wgs84/field-surveys.geojson', encoding='utf-8') as fh:
    _AREA_STATS = json.load(fh)['areaStats']


def area_row(key):
    """(spaces removed, displaced as % of removed, on-street absorption %) for an area."""
    s = {x['label']: x['value'] for x in _AREA_STATS[key]['displacement']}
    removed = s['Spaces Removed']
    displaced = s['Cars Displaced (Peak Hour)']
    return removed, round(100.0 * displaced / removed) if removed else 0, s['% Absorbed On-Street']


# slide-4 row order → the shape names holding removed / displaced / re-absorbed
SLIDE4_ROWS = {
    'kentron': ('TextBox 21', 'TextBox 22', 'TextBox 23'),
    'mega':    ('TextBox 31', 'TextBox 32', 'TextBox 33'),
    'komitas': ('TextBox 42', 'TextBox 43', 'TextBox 44'),
    'garegin': ('TextBox 52', 'TextBox 53', 'TextBox 54'),
    'shiraz':  ('TextBox 63', 'TextBox 64', 'TextBox 65'),
    'malatia': ('TextBox 73', 'TextBox 74', 'TextBox 75'),
}

SRC = r'C:/Users/user/Yerevan-Parking/Final Presentation/Parking Slides - 22072026.pptx'
DST = r'C:/Users/user/Yerevan-Parking/Final Presentation/Parking Slides - 23072026.pptx'

# --- figures (C1+C2 only) -----------------------------------------------------------
ON_CORRIDOR = 6095          # impact=corridor, C1+C2
ON_STREET, ON_SEGS = 7825, 606      # corridor + buffer, matching the old 8,862 basis
OFF_STREET, OFF_FAC = 6732, 239     # excludes the withdrawn NalbandyanPavstocByuzand
TOTAL = ON_STREET + OFF_STREET      # 14,557
# per-corridor removal, derived rather than typed: this line drifted out of step with
# the total once Nalbandyan016's 14 spaces left the Corridor 1 base (2,363 -> 2,349),
# leaving the deck showing 1,494 + 3,395 = 4,889 beside a card reading 4,875.
C1_EXISTING, C1_RETAINED = 2349, 869
C2_EXISTING, C2_RETAINED = 3746, 351
C1_REMOVED = C1_EXISTING - C1_RETAINED      # 1,480
C2_REMOVED = C2_EXISTING - C2_RETAINED      # 3,395
REMOVED, RETAINED = C1_REMOVED + C2_REMOVED, C1_RETAINED + C2_RETAINED
assert (REMOVED, RETAINED) == (4875, 1220), (REMOVED, RETAINED)
FREE, ZONE_B, ZONE_A, TAXI = 5248, 534, 246, 67
assert FREE + ZONE_B + ZONE_A + TAXI == ON_CORRIDOR == C1_EXISTING + C2_EXISTING

pc = lambda v: 100.0 * v / ON_CORRIDOR

TEXT_EDITS = {
    2: {
        'TextBox 3':  'Across Corridors 1 and 2 and 100 m buffer',
        'TextBox 8':  f'{TOTAL:,}',
        'TextBox 12': f'{ON_STREET:,}',
        'TextBox 13': f'On-street, across {ON_SEGS} segments',
        'TextBox 16': f'{ON_CORRIDOR:,}',
        'TextBox 20': f'{OFF_STREET:,}',
        'TextBox 21': f'Off-street, across {OFF_FAC} facilities',
        'TextBox 22': f'On-corridor on-street supply by regulatory zone ({ON_CORRIDOR:,} spaces)',
        'TextBox 28': f'Free / unregulated  {FREE:,} ({pc(FREE):.1f}%)',
        'TextBox 30': f'Zone B (blue) – paid  {ZONE_B} ({pc(ZONE_B):.1f}%)',
        'TextBox 32': f'Zone A (red) – paid  {ZONE_A} ({pc(ZONE_A):.1f}%)',
        'TextBox 34': f'Taxi  {TAXI} ({pc(TAXI):.1f}%)',
        'TextBox 37': f'{pc(FREE):.0f}% of corridor parking is free of charge — and barely organized',
        # markings rises 12% -> 14% because every marked space is on C1+C2; Nalbandyan
        # drops 83 -> 69 with the out-of-boundary segment gone.
        'TextBox 38': ('•  Only 20% has signage and 14% has markings. Zone A is concentrated in '
                       'Kentron (Nalbandyan 69, Abovyan 46, Amiryan 43); Zone B is almost entirely '
                       'Komitas Avenue (522 of 534).'),
    },
    4: {
        # These two columns are both percentages with DIFFERENT bases — displaced is over
        # spaces removed, re-absorbed is over displaced cars — so each heading names its
        # own denominator. Read as a matching pair they invite the misreading that both
        # are shares of demand.
        'TextBox 10': 'Displaced\n(% of spaces removed)',
        'TextBox 11': 'Re-absorbed\n(% of displaced)',
        'TextBox 80': (f'•  Kentron and Komitas peak over capacity yet re-absorb only '
                       f'{area_row("kentron")[2]}% and {area_row("komitas")[2]}% nearby; Malatia-Sebastia '
                       f'peaks at 54% and re-absorbs {area_row("malatia")[2]}% — so the package is applied '
                       'by zone, not uniformly.'),
        # The two columns rest on different time bases and a reader doing the arithmetic
        # will notice; state it rather than let it be discovered.
        'TextBox 81': ('Peak occupancy = distinct vehicles over the busiest hour ÷ spaces; above 100% means '
                       'more distinct vehicles than spaces used the kerb in that hour. It sums each zone\'s '
                       'own peak hour, so zones busiest at different times are added together. Displaced '
                       '(% of spaces removed) instead counts one shared clock hour — a stricter basis, so '
                       'the two columns are not directly comparable; on the single-hour basis Kentron reads '
                       '89% and Komitas 83%. Surveyed areas only, not corridor-wide. Source: field occupancy '
                       'survey, 29 May – 3 June 2026.'),
        **{name: value for key, shapes in SLIDE4_ROWS.items()
           for name, value in zip(shapes, (f'{area_row(key)[0]:,}', f'{area_row(key)[1]}%', f'{area_row(key)[2]}%'))},
    },
    6: {
        'TextBox 3':  f'The conceptual cross-section removes {100 * REMOVED / ON_CORRIDOR:.0f}% of on-corridor parking',
        'TextBox 6':  'CORRIDORS 1 & 2 · impact of the conceptual design',
        'TextBox 7':  f'Corridor 1: {C1_REMOVED:,}  ·  Corridor 2: {C2_REMOVED:,}',
        'TextBox 10': f'{ON_CORRIDOR:,}',
        'TextBox 14': f'{REMOVED:,}',
        'TextBox 18': f'{100 * REMOVED / ON_CORRIDOR:.0f}%',
        'TextBox 22': f'{RETAINED:,}',
        'TextBox 30': f'{D_RATIO}%',
        'TextBox 34': f'{D_ABSORB}%',
        'TextBox 38': f'{YARD_DEP}%',
        'TextBox 40': (f'Open and manage the yards BEFORE removing kerb: nearby streets absorb only '
                       f'{D_ABSORB}% (10% in Kentron, 0% in Komitas). Re-absorption reaches ~100% only '
                       f'via off-street, and {YARD_DEP}% of that is gated residential yards.'),
        'TextBox 42': (f'The {D_RATIO}% is the displaced demand set against the spaces removed in the same six '
                       f'areas ({D_DISPLACED:,} vehicles ÷ {D_REMOVED:,} spaces): the kerb is not full at the '
                       'moment, so fewer vehicles need re-homing than spaces are lost. Corridor-wide figures '
                       'cover Corridors 1 and 2 only — Corridor 3 is excluded pending its conceptual design '
                       '— and will be updated once that design is finalized.'),
    },
}

# The left-hand chart on slide 6 renders the removed/displaced pair as an image, so the
# text edits above would otherwise leave it showing the superseded 1,643 / 908.
PICTURE_SWAPS = {6: {'Picture 27': 'Final Presentation/.slide-assets/chart_displacement.png'}}

# Stacked regulatory bar on slide 2: the segment widths encode the shares, so the text
# edits above are not enough on their own. Rebuilt left-to-right across the same span.
BAR_L, BAR_W = 0.240, 9.520
BAR_SEGMENTS = [('Rectangle 23', FREE), ('Rectangle 24', ZONE_B),
                ('Rectangle 25', ZONE_A), ('Rectangle 26', TAXI)]


def set_text(shape, new):
    """Replace a shape's text, keeping each paragraph's existing run formatting.

    Hand-editing in PowerPoint has split these paragraphs into several runs that all
    share the same size/bold/colour, so collapsing to one run per paragraph is lossless
    here — and far safer than text_frame.text, which drops formatting entirely.

    Newlines map onto existing paragraphs one-for-one (the two-line column headings on
    slide 4 are genuine separate paragraphs, not line breaks inside one). Surplus
    paragraphs are removed; asking for more lines than the shape has is an error rather
    than a silent half-write.
    """
    paras = shape.text_frame.paragraphs
    lines = new.split('\n')
    if len(lines) > len(paras):
        raise ValueError(f'{shape.name}: {len(lines)} lines requested but shape has {len(paras)} paragraphs')
    for para, line in zip(paras, lines):
        if not para.runs:
            raise ValueError(f'{shape.name}: no runs to inherit formatting from')
        para.runs[0].text = line
        for run in para.runs[1:]:
            run._r.getparent().remove(run._r)
    for para in paras[len(lines):]:
        para._p.getparent().remove(para._p)


def replace_picture(slide, shape, image_path):
    """Swap a picture's image, keeping its exact position, size and z-order.

    Inserting and deleting is safer than rewriting the image part in place: image parts
    are shared between shapes that use the same file, so overwriting the blob would hit
    every other slide using it. The new element is moved back to the old one's index so
    stacking order is preserved.
    """
    old = shape._element
    idx = list(old.getparent()).index(old)
    new_shape = slide.shapes.add_picture(image_path, shape.left, shape.top,
                                         shape.width, shape.height)
    new = new_shape._element
    new.getparent().remove(new)
    old.getparent().insert(idx, new)
    old.getparent().remove(old)
    return new_shape


shutil.copyfile(SRC, DST)
prs = Presentation(DST)

applied, missing = 0, []
for idx, edits in TEXT_EDITS.items():
    slide = prs.slides[idx - 1]
    by_name = {sh.name: sh for sh in slide.shapes}
    for name, new in edits.items():
        if name not in by_name:
            missing.append(f'slide {idx} / {name}')
            continue
        before = by_name[name].text_frame.text
        set_text(by_name[name], new)
        if before != new:
            print(f'  s{idx} {name:<12} {before[:58]!r}\n       -> {new[:58]!r}')
        applied += 1

for idx, swaps in PICTURE_SWAPS.items():
    slide = prs.slides[idx - 1]
    by_name = {sh.name: sh for sh in slide.shapes}
    for name, path in swaps.items():
        if name not in by_name:
            missing.append(f'slide {idx} / {name}')
            continue
        replace_picture(slide, by_name[name], path)
        print(f'  s{idx} {name:<12} image <- {path}')

# rebuild the stacked bar
slide2 = {sh.name: sh for sh in prs.slides[1].shapes}
cursor = BAR_L
for name, value in BAR_SEGMENTS:
    if name not in slide2:
        missing.append(f'slide 2 / {name}')
        continue
    width = BAR_W * value / ON_CORRIDOR
    slide2[name].left = Inches(cursor)
    slide2[name].width = Inches(width)
    print(f'  s2 {name:<12} L={cursor:.3f}" W={width:.3f}"  ({value} = {pc(value):.1f}%)')
    cursor += width
assert abs(cursor - (BAR_L + BAR_W)) < 1e-6, f'bar segments end at {cursor}, expected {BAR_L + BAR_W}'

prs.save(DST)
print(f'\n{applied} text edits applied; bar rebuilt.')
if missing:
    print('MISSING SHAPES (nothing written for these):')
    for m in missing:
        print('   ', m)
print('Saved ->', DST)
