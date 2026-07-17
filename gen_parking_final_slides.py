# -*- coding: utf-8 -*-
"""Parking slides for the final ADB mission presentation.

Builds two decks from one shared data block, so the tight and full versions can
never contradict each other:

    Parking Slides - FULL (11).pptx    menu to pick from
    Parking Slides - TIGHT (5).pptx    the likely presentation cut

Both are constructed on "Final Presentation/ADB Mission_Template.pptx" so they
carry the ADB master (olive top bar, footer band, logos) and can be pasted
straight into the master mission deck.

Numbers come from the two final .docx deliverables (Output 8 / Output 10) and
from app/static/data/wgs84/field-surveys.geojson, which is the single source of
truth per CLAUDE.md. See DATA below for the provenance of each figure.

Visual language follows "Parking Assessment - 30042026.pptx" (stat tiles, thin
grey cards, coloured sensitivity triad). Fonts are deliberately NOT Roboto Slab:
that deck references fonts absent from every machine checked (including the ADB
theme's own "Ideal Sans Semibold"), so it already renders substituted. Calibri /
Arial render identically everywhere.
"""

import copy
import os
import zipfile
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import gen_parking_slide_charts as charts

BASE = "C:/Users/user/Yerevan-Parking/Final Presentation"
TEMPLATE = BASE + "/ADB Mission_Template.pptx"
OUT_FULL = BASE + "/Parking Slides - FULL (11).pptx"
OUT_TIGHT = BASE + "/Parking Slides - TIGHT (5).pptx"
ASSETS = BASE + "/.slide-assets"
O8_DOCX = BASE + "/Output 8 - Parking Surveys and Analysis Report 23062026 (rev).docx"

# Figure 2 of Output 8 — the app's "Parking Regulation" step. Verified to carry the
# CORRECTED split (6,144 / 260 / 534 / 67 = 7,005). Note the LIVE app still renders
# 583 for Zone B, so re-shoot it from the app only after the app data is fixed.
# Figure 1 (image3, total 16,341) and Figure 7 (image9, off-street 7,479) are STALE
# — do not use them.
DOCX_FIGURES = {"app_regulation.png": "word/media/image4.png"}

# --- palette (sampled from Parking Assessment - 30042026.pptx) -----------------
NAVY = RGBColor(0x00, 0x20, 0x60)   # ADB template title colour
INK = RGBColor(0x1F, 0x1F, 0x1F)
GREY = RGBColor(0x66, 0x66, 0x66)
CARD_BG = RGBColor(0xFA, 0xFA, 0xFA)
CARD_LINE = RGBColor(0xE0, 0xE0, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC6, 0x28, 0x28)    # Zone A
BLUE = RGBColor(0x15, 0x65, 0xC0)   # Zone B
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xD4, 0xA5, 0x20)  # taxi
FREE_GR = RGBColor(0xE8, 0xE8, 0xE8)
BAR_GR = RGBColor(0xED, 0xED, 0xED)

HIGH_T, HIGH_BG = RGBColor(0x79, 0x1F, 0x1F), RGBColor(0xFC, 0xEB, 0xEB)
MED_T, MED_BG = RGBColor(0x63, 0x38, 0x06), RGBColor(0xFA, 0xEE, 0xDA)
LOW_T, LOW_BG = RGBColor(0x08, 0x50, 0x41), RGBColor(0xE1, 0xF5, 0xEE)
BLUE_BG = RGBColor(0xE3, 0xF0, 0xFB)

FONT = "Calibri"
TAB = "Parking Assessment"

# geometry: April deck's tile grid, dropped to y=1.55 to clear the layout logo
# that sits at (8.95, 0.30)-(10.00, 1.49) on the Blank layout.
L, R = 0.24, 9.76
CW = R - L                      # 9.52 content width
ROW1 = 1.55
TILE_H = 1.16


def grid(n, gap=0.20):
    """x positions and width for n equal columns across the content width."""
    w = (CW - gap * (n - 1)) / n
    return [L + i * (w + gap) for i in range(n)], w


# --- data ---------------------------------------------------------------------
# Provenance: O8 = Output 8 (rev).docx, O10 = Output 10 (rev).docx,
# GJ = field-surveys.geojson areaStats/displacement.
DATA = dict(
    # supply (O8 Table 1 / Aggregate Supply Summary) -- CORRECTED totals.
    # The April deck's 16,341 / 7,479 are superseded; do not reintroduce.
    total=15897, onstreet=8862, offstreet=7035, segments=728, facilities=266,
    on_corridor=7005, cross=1857,
    zone_a=260, zone_b=534, taxi=67, free=6144,
    signage_pct=20, marking_pct=12,
    # corridor-wide impact (O8) -- distinct scope from the surveyed-area figures
    removed=5953, retained=1052, removed_pct=85,
    c1_removed=1836, c2_removed=3310, c3_removed=807,
    # demand (O8 / GJ areaStats.all)
    peak_occ=93, over85=69, over_cap_zones=116,
    d_1h=63, d_2h=77, d_worker=6, d_allday=7, avg_stay=1.8,
    # displacement (GJ displacement) -- SURVEYED AREAS ONLY
    d_removed=1643, d_displaced=908, d_ratio=55,
    absorb_onstreet=43, absorb_total=100, yard_dep=92,
)

# peak occupancy + on-street absorption by area (GJ areaStats)
AREAS = [
    ("Kentron", 138, 14),
    ("Komitas", 123, 0),
    ("Gai Avenue (Mega Mall)", 123, 68),
    ("Garegin Nzhdeh", 92, 30),
    ("Shiraz / Hasratyan", 75, 69),
    ("Malatia-Sebastia", 54, 97),
]

# surveyed off-street facilities (O8, Available Absorptive Capacity table)
OFFSTREET = [
    ("Gai Avenue (Mega Mall)", "Palace lot", 32, 115, 194),
    ("Kentron", "Nalbandyan yard", 60, 81, 142),
    ("Komitas", "Komitas City lot", 123, 56, 101),
    ("Malatia-Sebastia", "Sebastia yard", 101, 32, 77),
    ("Shiraz / Hasratyan", "Shiraz off-street", 71, 37, 59),
    ("Garegin Nzhdeh", "GN off-street", 40, 22, 45),
]

# measure verdicts (O10 "Field-survey verdict" paragraphs)
MEASURES = [
    ("Extend paid zones", "CONFIRMED — LEAD", GREEN, LOW_BG,
     "About 90% of the kerb is unpriced. Roll the red/blue-line regime onto side "
     "streets within 100 m, with ANPR enforcement."),
    ("Organise free parking", "CONFIRMED — WIDEN", GREEN, LOW_BG,
     "88% of corridor parking is unmarked and 24.5% sits off the carriageway. "
     "Widen beyond outer zones: Komitas (35% informal), Kentron (60% angled)."),
    ("Open residential yards", "NOW A PRECONDITION", HIGH_T, HIGH_BG,
     "Carries about 92% of the absorptive capacity. Re-absorption is conditional "
     "on opening these yards — no longer an optional outer-zone pilot."),
    ("Delivery spaces", "CONFIRMED — RE-SIZE", MED_T, MED_BG,
     "Size loading bays to the measured freight pockets — truck share 12.6% in "
     "Malatia-Sebastia and 8.0% in Shiraz against 4.9% citywide — not uniformly."),
    ("Resident permits", "KEEP — BUT PRICE IT", BLUE, BLUE_BG,
     "Tbilisi is a cautionary case, not a model: two free vehicles per apartment "
     "produced exactly the long-stay pattern to avoid. Price or cap from the outset."),
    ("Visitor caps", "RE-SCOPE", MED_T, MED_BG,
     "The data contradicts the premise: 63% of vehicles already stay an hour or "
     "less and all-day parking is ~6%. Demote to a targeted contingency; lead with pricing."),
    ("Park & ride", "FUTURE STUDY", GREY, BAR_GR,
     "The daytime fleet is overwhelmingly short-stay, so P&R is not the binding "
     "lever now. Retain as a strategic study item for the BRT terminals."),
]

ZONES = [
    ("HIGH SENSITIVITY", HIGH_T, HIGH_BG,
     "Kentron CBD · 138% peak · 14% absorbed\nKomitas · 123% peak · 0% absorbed\n"
     "— moved up from medium by the survey",
     ["Extend paid zones — lead", "Open residential yards — precondition",
      "Delivery spaces", "Resident permits", "Visitor caps — contingency only"]),
    ("MEDIUM SENSITIVITY", MED_T, MED_BG,
     "Arshakunyats · Kievyan\nSubstantial pressure, but more nearby capacity "
     "than the core\n— orderly redistribution",
     ["Extend paid zones", "Delivery spaces", "Park & ride where relevant"]),
    ("LOWER SENSITIVITY", LOW_T, LOW_BG,
     "Shiraz 75% · Malatia-Sebastia 54% peak\nBagratunyats, Sebastia, Raffi, "
     "Nor Nork\n— structure before pricing",
     ["Organise free parking", "Park & ride", "Open residential yards"]),
]


# --- assets -------------------------------------------------------------------
def ensure_assets():
    """Charts + the app screenshot lifted straight out of the .docx deliverable.

    Extracting from the .docx rather than committing a copy keeps one source: if
    the report figure is regenerated, so is the slide.
    """
    os.makedirs(ASSETS, exist_ok=True)
    charts.build()
    with zipfile.ZipFile(O8_DOCX) as z:
        for out, member in DOCX_FIGURES.items():
            with open(os.path.join(ASSETS, out), "wb") as fh:
                fh.write(z.read(member))
    print("assets ->", ASSETS)


def add_img(slide, name, x, y, w=None, h=None):
    """Place an image preserving its aspect ratio. Give w OR h."""
    path = os.path.join(ASSETS, name)
    iw, ih = Image.open(path).size
    ar = ih / float(iw)
    if w is None:
        w = h / ar
    else:
        h = w * ar
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    Inches(w), Inches(h))


# --- primitives ---------------------------------------------------------------
def new_deck():
    prs = Presentation(TEMPLATE)
    for i in range(len(prs.slides._sldIdLst) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    return prs


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, lines, size=9, color=GREY, bold=False, font=FONT, space=2, italic=False):
    """lines: str or list of str (or (text, kwargs) tuples) -> paragraphs."""
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        over = {}
        if isinstance(line, tuple):
            line, over = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = tf.paragraphs[0].alignment
        p.space_after = Pt(over.get("space", space))
        r = p.add_run()
        r.text = line
        f = r.font
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.name = over.get("font", font)
        f.color.rgb = over.get("color", color)
    return tf


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.text_frame.word_wrap = True
    return s


def slide_number(slide):
    tf = textbox(slide, 8.99, 5.28, 1.01, 0.30, align=PP_ALIGN.RIGHT)
    p = tf.paragraphs[0]
    fld = p._p.makeelement(qn("a:fld"), {"id": "{B6F15528-21DE-4FA1-8B65-3B0972ABAC2C}",
                                         "type": "slidenum"})
    rPr = p._p.makeelement(qn("a:rPr"), {"lang": "en-US", "sz": "1000"})
    latin = p._p.makeelement(qn("a:latin"), {"typeface": FONT})
    rPr.append(latin)
    t = p._p.makeelement(qn("a:t"), {})
    t.text = "2"
    fld.append(rPr)
    fld.append(t)
    p._p.append(fld)


def base_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # Blank
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    # section tab: sits on the master's olive bar; 11pt italic, theme major font
    tf = textbox(slide, 0.0, -0.008, 5.03, 0.286, anchor=MSO_ANCHOR.MIDDLE)
    tf.margin_left = Inches(0.14)
    r = tf.paragraphs[0].add_run()
    r.text = TAB
    r.font.size, r.font.italic, r.font.color.rgb = Pt(11), True, INK
    r.font._rPr.append(r.font._rPr.makeelement(qn("a:latin"), {"typeface": "+mj-lt"}))
    write(textbox(slide, L, 0.42, 8.60, 0.34), title, size=18, bold=True, color=NAVY)
    write(textbox(slide, L, 0.80, 8.60, 0.34), subtitle, size=11.5, color=GREY)
    rect(slide, L, 1.22, 8.61, 0.014, CARD_LINE)
    slide_number(slide)
    return slide


def tile(slide, x, y, w, value, label, color=INK, h=TILE_H, accent=None, vsize=20):
    rect(slide, x, y, w, h, CARD_BG, CARD_LINE)
    if accent:
        rect(slide, x, y, w, 0.045, accent)
    write(textbox(slide, x, y + (0.16 if h >= 1.0 else 0.12), w, 0.46,
                  align=PP_ALIGN.CENTER), value, size=vsize, bold=True, color=color)
    write(textbox(slide, x + 0.08, y + (0.66 if h >= 1.0 else 0.44), w - 0.16, 0.38,
                  align=PP_ALIGN.CENTER), label, size=8.5, color=GREY, space=0)


def footnote(slide, text, y=4.86):
    write(textbox(slide, L, y, CW, 0.30), text, size=8, color=GREY, italic=True)


def card(slide, x, y, w, h, heading, body, accent=None, bg=CARD_BG,
         line=CARD_LINE, hcolor=INK, bullets=None, bsize=8.5, bspace=4):
    rect(slide, x, y, w, h, bg, line)
    cy = y + 0.10
    if accent:
        rect(slide, x, y, w, 0.045, accent)
        cy = y + 0.16
    write(textbox(slide, x + 0.12, cy, w - 0.24, 0.24), heading,
          size=10, bold=True, color=hcolor)
    cy += 0.30
    if body:
        tf = textbox(slide, x + 0.12, cy, w - 0.24, 0.5)
        write(tf, body.split("\n"), size=8.5, color=GREY, space=3)
        cy += 0.20 + 0.16 * (body.count("\n") + 1)
    if bullets:
        tf = textbox(slide, x + 0.12, cy, w - 0.24, max(0.2, y + h - cy - 0.10))
        write(tf, ["•  " + b for b in bullets], size=bsize, color=INK, space=bspace)


def hbar(slide, x, y, w, h, pct, scale, color, label, value, ref=None):
    """One horizontal bar row: name | bar | value."""
    write(textbox(slide, x, y + 0.06, 1.62, 0.26, align=PP_ALIGN.RIGHT),
          label, size=8.5, color=INK)
    bx = x + 1.72
    rect(slide, bx, y, w, h, BAR_GR)
    rect(slide, bx, y, min(pct, 100 if ref is None else pct) * scale, h, color)
    write(textbox(slide, bx + pct * scale + 0.08, y + 0.06, 0.7, 0.26),
          value, size=8.5, bold=True, color=color)


def stacked_zone_bar(slide, y):
    """On-corridor supply by regulatory zone (7,005 spaces)."""
    write(textbox(slide, L, y, CW, 0.22),
          "On-corridor on-street supply by regulatory zone (7,005 spaces)",
          size=9, bold=True, color=INK)
    segs = [("Free / unregulated", DATA["free"], FREE_GR, INK),
            ("Zone B (blue) – paid", DATA["zone_b"], BLUE, BLUE),
            ("Zone A (red) – paid", DATA["zone_a"], RED, RED),
            ("Taxi", DATA["taxi"], AMBER, AMBER)]
    x, by = L, y + 0.28
    for _, val, col, _ in segs:
        w = CW * val / DATA["on_corridor"]
        rect(slide, x, by, w, 0.30, col)
        x += w
    lx = L
    for name, val, col, tcol in segs:
        rect(slide, lx, by + 0.44, 0.11, 0.11, col)
        write(textbox(slide, lx + 0.17, by + 0.41, 2.05, 0.20),
              "%s  %s (%.1f%%)" % (name, "{:,}".format(val),
                                   100.0 * val / DATA["on_corridor"]),
              size=8, color=GREY)
        lx += 2.38


# --- slides -------------------------------------------------------------------
def s_whats_new(prs):
    s = base_slide(prs, "Parking: what has changed since April",
                   "The demand side is now measured, not inferred — six representative "
                   "areas surveyed 29 May – 3 June 2026")
    xs, w = grid(4)
    tiles = [("93%", "Peak occupancy, capacity-weighted", GREEN),
             ("63%", "Of vehicles stay one hour or less", BLUE),
             ("55%", "Displaced demand vs. spaces removed", RED),
             ("5–11", "Vehicles per space per day (turnover)", INK)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    xs2, w2 = grid(2)
    card(s, xs2[0], 3.05, w2, 1.42, "April 2026  —  inferred", None, accent=GREY,
         bsize=9.5, bspace=7,
         bullets=["Occupancy survey descoped on proportionality grounds",
                  "Displacement argued from European BHLS precedent — Nantes, "
                  "Metz, Rouen, Eindhoven, Helsingborg",
                  "Re-absorption asserted, not quantified"])
    card(s, xs2[1], 3.05, w2, 1.42, "June 2026  —  measured in Yerevan", None,
         accent=GREEN, bsize=9.5, bspace=7,
         bullets=["908 cars displaced at the peak hour against 1,643 spaces removed",
                  "The European finding is now confirmed by local measurement",
                  "Re-absorption quantified — and its precondition exposed"])
    footnote(s, "Six areas: Kentron, Komitas, Gai Avenue, Garegin Nzhdeh, "
                "Shiraz/Hasratyan, Malatia-Sebastia. Hourly licence-plate sweeps, "
                "07:00–24:00, on four ordinary mid-week working days.")
    return s


def s_approach(prs):
    s = base_slide(prs, "Approach: how the evidence was built",
                   "A desk-based supply inventory, completed by a targeted field "
                   "occupancy survey")
    xs, w = grid(2)
    card(s, xs[0], ROW1, w, 2.15, "1  ·  Desk-based supply inventory", None,
         accent=BLUE, bsize=9.5, bspace=8,
         bullets=["994 geo-referenced features — 728 on-street segments and "
                  "266 off-street facilities",
                  "Built from 360° street video (yerevan.omni7.app), satellite "
                  "and street-level imagery",
                  "Parallel capacity = length ÷ 7.5 m; angled bays counted "
                  "individually",
                  "~15% of spaces independently cross-checked, plus on-the-spot "
                  "field visits",
                  "Published live at yerevan-parking.vercel.app"])
    card(s, xs[1], ROW1, w, 2.15, "2  ·  Field occupancy survey (targeted)", None,
         accent=GREEN, bsize=9.5, bspace=8,
         bullets=["Six representative areas spanning the corridor sensitivity range",
                  "Hourly licence-plate sweeps, 07:00–24:00, four ordinary "
                  "working days",
                  "Records time, zone, plate, vehicle type, kerb location, parking "
                  "method and legality",
                  "Yields occupancy, duration of stay and turnover directly",
                  "Areas spread across separate days, so a one-off disruption shows "
                  "as a local outlier"])
    card(s, L, 4.02, CW, 0.74,
         "The ToR is met in substance, not departed from", None, accent=NAVY, bsize=9,
         bullets=["The occupancy component, initially descoped on proportionality "
                  "grounds, is now satisfied in representative form — the "
                  "deliverable meets the intent of the ToR rather than departing from it."])
    return s


def s_supply(prs):
    s = base_slide(prs, "Supply baseline: 15,897 spaces surveyed",
                   "Three bus priority corridors (34.6 km) and their 100-metre "
                   "influence area")
    xs, w = grid(4)
    tiles = [("15,897", "Total spaces surveyed", INK),
             ("8,862", "On-street spaces (728 segments)", BLUE),
             ("7,035", "Off-street spaces (266 facilities)", RGBColor(0x7C, 0x4D, 0xFF)),
             ("7,005", "On-street, on the corridors themselves", INK)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c)
    stacked_zone_bar(s, 2.92)
    card(s, L, 4.02, CW, 0.72,
         "88% of corridor parking is free — and barely organised", None,
         accent=AMBER,
         bullets=["Only 20% of on-corridor spaces have signage and only 12% have "
                  "markings. The supply is not just cheap, it is unstructured — "
                  "which is why organising it is a measure in its own right."])
    footnote(s, "Totals corrected June 2026 (off-street 7,479 → 7,035; total "
                "16,341 → 15,897). On-street and cross-street figures are unchanged.")
    return s


def s_occupancy(prs):
    s = base_slide(prs, "Demand: the kerb is saturated at the peak",
                   "Capacity-weighted peak occupancy of 93%, with 69% of zones above "
                   "the 85% efficiency threshold")
    scale = 0.0295          # inches per percentage point
    y = ROW1 + 0.14
    write(textbox(s, L, ROW1 - 0.16, 4.0, 0.20),
          "Peak occupancy by surveyed area", size=9, bold=True, color=INK)
    # 85% reference line, drawn behind the bars
    refx = L + 1.72 + 85 * scale
    rect(s, refx, ROW1 + 0.06, 0.01, 2.62, RGBColor(0x99, 0x99, 0x99))
    write(textbox(s, refx - 0.42, ROW1 + 2.70, 0.9, 0.18), "85% threshold",
          size=7, color=GREY, italic=True)
    for name, pct, _ in AREAS:
        col = RED if pct >= 100 else (AMBER if pct >= 85 else GREEN)
        hbar(s, L, y, 100 * scale, 0.30, pct, scale, col, name, "%d%%" % pct, ref=85)
        y += 0.42
    xs, w = grid(1)
    bx = 6.70
    for i, (v, l, c) in enumerate([("93%", "Peak occupancy, capacity-weighted", GREEN),
                                   ("69%", "Of zones above the 85% threshold", AMBER),
                                   ("116", "Zones over legal capacity at the peak", RED)]):
        tile(s, bx, ROW1 + i * 0.98, 3.06, v, l, color=c, h=0.86, vsize=17)
    footnote(s, "Occupancy above 100% is chronic overflow — more vehicles present "
                "than legal striped capacity. Only ~0.2% of vehicles were flagged as "
                "illegally parked: the issue is saturation and tolerated informality, "
                "not overt illegality.")
    return s


def s_duration(prs):
    s = base_slide(prs, "Demand: a short-stay, high-turnover kerb",
                   "Parking is dominated by short visits; genuine all-day storage is "
                   "a small tail")
    xs, w = grid(4)
    tiles = [("63%", "Stay one hour or less", BLUE),
             ("77%", "Stay two hours or less", BLUE),
             ("1.8 h", "Average stay", INK),
             ("5–11", "Vehicles per space per day", GREEN)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c)
    xs3, w3 = grid(3)
    groups = [("Visitors & customers  —  ~63%", BLUE,
               "The dominant group: short retail and errand access. This is the "
               "demand the kerb actually serves."),
              ("Commuters & workers  —  ~6%", GREEN,
               "A small all-day tail. Better addressed as a BRT mode-shift "
               "opportunity than as a fining target."),
              ("Residents  —  ~7% overnight", MED_T,
               "Under-counted by a daytime window (07:00–24:00). To be "
               "protected by the permit scheme, not squeezed.")]
    for x, (h, c, b) in zip(xs3, groups):
        card(s, x, 3.05, w3, 1.22, h, b, accent=c, hcolor=c)
    footnote(s, "Implication: visitor caps address a problem the data does not show. "
                "Turnover is already high, so pricing — not maximum-stay "
                "enforcement — is the binding lever.", y=4.55)
    return s


def s_impact(prs):
    s = base_slide(prs, "Impact: 85% of on-corridor parking is removed",
                   "Inherent to the central-median cross-section — a design "
                   "consequence, not an incidental one")
    xs, w = grid(3)
    tiles = [("5,953", "On-street spaces removed", RED),
             ("1,052", "Spaces retained or re-established", GREEN),
             ("85%", "Of the on-corridor supply removed", RED)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    cs = [("Corridor 1  —  11.66 km", "1,836 spaces removed\nBagratunyats "
           "(south) to Kentron; carries the Zone A core."),
          ("Corridor 2  —  19.19 km", "3,310 spaces removed\nNor Nork (east) to "
           "Davtashen; the largest single loss."),
          ("Corridor 3  —  3.81 km", "807 spaces removed\nSebastia/Malatia; "
           "entirely residential off-street context.")]
    for x, (h, b) in zip(xs, cs):
        card(s, x, 3.00, w, 1.42, h, b, accent=NAVY)
    footnote(s, "Cross-section: 2 × 3.5 m bus lanes, 2 × 3.0 m mixed traffic "
                "per direction, 0.5 m separators, cycle lanes and 2.50 m sidewalks — "
                "leaving no carriageway width for parking. Consistent with the Mott "
                "MacDonald (2022) strategy and the Yerevan Sustainable Urban Transport Strategy.")
    return s


def s_displacement(prs):
    s = base_slide(prs, "Displacement: measured, not inferred",
                   "Peak displaced demand equals only 55% of the spaces removed — "
                   "the European finding, now measured in Yerevan")
    xs, w = grid(3)
    tiles = [("1,643", "Spaces removed in the surveyed areas", GREY),
             ("908", "Cars actually displaced at the peak hour", RED),
             ("55%", "Displaced demand as a share of supply removed", GREEN)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    write(textbox(s, L, 2.92, 5.0, 0.20),
          "Share of displaced demand absorbed on nearby streets",
          size=9, bold=True, color=INK)
    y, scale = 3.16, 0.030
    for name, _, ab in AREAS:
        col = GREEN if ab >= 60 else (AMBER if ab >= 25 else RED)
        hbar(s, L, y, 100 * scale, 0.21, ab, scale, col, name, "%d%%" % ab)
        y += 0.26
    card(s, 6.70, 3.16, 3.06, 1.44, "Why this matters", None, accent=GREEN,
         bullets=["Not every removed space needs replacing one-for-one",
                  "Some users shift mode, retime, or were parked informally",
                  "But on-street alone absorbs just 43% — the rest must come "
                  "from off-street"])
    footnote(s, "Surveyed areas only — a subset of the corridor-wide 5,953 "
                "removals. Measured at the peak hour, the binding case.", y=4.80)
    return s


def s_yards(prs):
    s = base_slide(prs, "The catch: re-absorption depends on opening gated yards",
                   "About 92% of the absorptive capacity is gated residential yards, "
                   "counted at gross capacity as if already open")
    xs, w = grid(3)
    tiles = [("~100%", "Aggregate re-absorption — best case", GREEN),
             ("43%", "Absorbed by surviving on-street capacity", AMBER),
             ("92%", "Of absorptive capacity sits behind gates", HIGH_T)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    # observed off-street occupancy table
    write(textbox(s, L, 2.90, 6.2, 0.20),
          "Observed occupancy in the six surveyed off-street facilities (427 spaces)",
          size=9, bold=True, color=INK)
    ty = 3.12
    hdr = [("Area", 0.0, 1.95), ("Facility", 1.95, 1.60), ("Spaces", 3.55, 0.62),
           ("Avg.", 4.17, 0.62), ("Peak", 4.79, 0.62)]
    rect(s, L, ty, 5.41, 0.22, BAR_GR)
    for name, dx, ww in hdr:
        al = PP_ALIGN.LEFT if dx < 3.5 else PP_ALIGN.RIGHT
        write(textbox(s, L + dx + 0.06, ty + 0.04, ww - 0.10, 0.18, align=al),
              name, size=7.5, bold=True, color=INK)
    ty += 0.22
    for area, fac, sp, avg, peak in OFFSTREET:
        over = peak >= 100
        rect(s, L, ty, 5.41, 0.22, WHITE if not over else RGBColor(0xFF, 0xF5, 0xF5),
             CARD_LINE)
        vals = [(area, 0.0, 1.95, PP_ALIGN.LEFT, INK, False),
                (fac, 1.95, 1.60, PP_ALIGN.LEFT, GREY, False),
                (str(sp), 3.55, 0.62, PP_ALIGN.RIGHT, GREY, False),
                ("%d%%" % avg, 4.17, 0.62, PP_ALIGN.RIGHT, GREY, False),
                ("%d%%" % peak, 4.79, 0.62, PP_ALIGN.RIGHT,
                 RED if over else GREEN, True)]
        for txt, dx, ww, al, col, bd in vals:
            write(textbox(s, L + dx + 0.06, ty + 0.04, ww - 0.10, 0.18, align=al),
                  txt, size=7.5, color=col, bold=bd)
        ty += 0.22
    write(textbox(s, L + 0.06, ty + 0.04, 4.2, 0.18),
          "All six — 427 spaces, ~52% full on average across the day",
          size=7.5, italic=True, color=GREY)
    card(s, 5.95, 3.12, 3.81, 1.54, "Why this is load-bearing", None, accent=HIGH_T,
         hcolor=HIGH_T,
         bullets=["Three of the six are already at or over capacity at the peak",
                  "Those facilities carry their own demand — excluded from the "
                  "absorptive count",
                  "The capacity that closes the gap is still behind gates",
                  "So: open and manage the yards BEFORE removing kerb"])
    footnote(s, "Kentron is ~86% yard-dependent, Komitas ~100%. Opening yards is "
                "largely untested in post-Soviet courtyards — ownership, consent and "
                "revenue-sharing need a dedicated study.", y=4.86)
    return s


def s_verdicts(prs):
    s = base_slide(prs, "What the measured evidence changed",
                   "Each mitigation measure re-tested against the field survey — "
                   "two were materially revised")
    xs, w = grid(4)
    for i, (name, chip, ccol, cbg, body) in enumerate(MEASURES):
        col, row = i % 4, i // 4
        x, y = xs[col], ROW1 + row * 1.66
        rect(s, x, y, w, 1.56, CARD_BG, CARD_LINE)
        rect(s, x, y, w, 0.045, ccol)
        write(textbox(s, x + 0.10, y + 0.14, w - 0.20, 0.20), name,
              size=9.5, bold=True, color=INK)
        chip_w = min(w - 0.20, 0.055 * len(chip) + 0.18)
        rect(s, x + 0.10, y + 0.40, chip_w, 0.20, cbg)
        write(textbox(s, x + 0.10, y + 0.435, chip_w, 0.16, align=PP_ALIGN.CENTER),
              chip, size=6.5, bold=True, color=ccol)
        write(textbox(s, x + 0.10, y + 0.68, w - 0.20, 0.80), body,
              size=7.5, color=GREY, space=0)
    card(s, xs[3], ROW1 + 1.66, w, 1.56,
         "Also flagged", None, accent=GREY, hcolor=GREY,
         bullets=["Phase out the flat annual permit — it removes the marginal "
                  "cost of dwell time and so works against the turnover the data "
                  "shows is the binding lever.",
                  "A candidate idea, not yet a measure: substantiate with PCS "
                  "permit-uptake and tariff data first."])
    footnote(s, "The corridor itself remains the primary mitigation: faster, more "
                "reliable public transport reduces the number of trips that need to "
                "park at all. Every measure above is complementary to that structural shift.")
    return s


def s_zones(prs):
    s = base_slide(prs, "Mitigation packages by sensitivity zone",
                   "Recalibrated by the survey: Komitas moves up to high sensitivity; "
                   "Gai Avenue is handled as an exception")
    xs, w = grid(3)
    for x, (name, tcol, bg, desc, measures) in zip(xs, ZONES):
        rect(s, x, ROW1, w, 2.42, bg, CARD_LINE)
        rect(s, x, ROW1, w, 0.045, tcol)
        write(textbox(s, x + 0.12, ROW1 + 0.14, w - 0.24, 0.20), name,
              size=10, bold=True, color=tcol)
        write(textbox(s, x + 0.12, ROW1 + 0.40, w - 0.24, 0.70), desc.split("\n"),
              size=8.5, color=INK, space=3)
        write(textbox(s, x + 0.12, ROW1 + 1.24, w - 0.24, 1.0),
              ["•  " + m for m in measures], size=8.5, color=tcol, space=5)
    card(s, L, 4.12, CW, 0.74,
         "Exception  —  Gai Avenue (Mega Mall)", None, accent=HIGH_T, hcolor=HIGH_T,
         bsize=9,
         bullets=["123% peak. Even filling every nearby yard leaves demand "
                  "unaccommodated (89% best case) — so the kerb parking stays for "
                  "now, to be revisited as the area changes."])
    return s


def s_conclusions(prs):
    s = base_slide(prs, "Conclusions and next steps",
                   "The parking case is now evidence-based end to end — with one "
                   "load-bearing condition")
    xs, w = grid(2)
    card(s, xs[0], ROW1, w, 2.45, "Conclusions", None, accent=NAVY,
         bullets=["Parking loss is substantial and inherent to the design — "
                  "about 85% of on-corridor supply (5,953 spaces).",
                  "The demand side is now measured, not inferred: 93% peak "
                  "occupancy, short-stay and high-turnover.",
                  "Displaced demand is only 55% of the spaces removed — the "
                  "European finding, confirmed locally.",
                  "Re-absorption approaching 100% is achievable but conditional on "
                  "opening gated yards (92% of capacity).",
                  "Mitigation must be layered, and must precede kerb removal rather "
                  "than follow it."])
    card(s, xs[1], ROW1, w, 2.45, "Next steps", None, accent=GREEN,
         bullets=["A dedicated residential-yard opening study: ownership, consent "
                  "thresholds, revenue-sharing, enforcement.",
                  "Extend the supply inventory citywide, beyond the corridor buffer.",
                  "Progressively expand the red-line zonal system into the "
                  "influence-area side streets.",
                  "Park & ride feasibility; resident-permit legal basis; "
                  "willingness-to-pay to calibrate tariffs.",
                  "Monitor side-street occupancy after commissioning and adapt "
                  "pricing and enforcement."])
    card(s, L, 4.12, CW, 0.68, "Sequencing is the decision that matters", None,
         accent=HIGH_T, hcolor=HIGH_T, bsize=9,
         bullets=["Removing kerb capacity ahead of opening the yards would strand "
                  "the displaced demand the survey measures. Mitigation before "
                  "removal is load-bearing, not advisory."])
    footnote(s, "Full interactive results: yerevan-parking.vercel.app  ·  Detail "
                "in Output 8 (Parking Surveys and Analysis) and Output 10 (Parking Analysis).")
    return s


# --- tight deck: visual-first --------------------------------------------------
# Text here is caption-only. Every slide leads with a chart or a map; the numbers
# live in tiles, not sentences.

def v_measured(prs):
    s = base_slide(prs, "Parking: the demand side is now measured",
                   "Six representative areas surveyed 29 May – 3 June 2026 — "
                   "occupancy, duration and turnover measured directly")
    xs, w = grid(4)
    tiles = [("93%", "Peak occupancy, capacity-weighted", GREEN),
             ("63%", "Stay one hour or less", BLUE),
             ("55%", "Displaced demand vs. spaces removed", RED),
             ("5–11", "Vehicles per space per day", INK)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    write(textbox(s, L, 2.86, 5.6, 0.20),
          "Vehicles parked, by hour of day — the kerb never empties between "
          "09:00 and 22:00", size=9, bold=True, color=INK)
    add_img(s, "chart_hourly.png", L, 3.06, w=5.18)
    card(s, 5.62, 3.06, 4.14, 1.94, "April  →  June", None, accent=GREEN, bsize=9,
         bspace=7,
         bullets=["Then: displacement inferred from European BHLS precedent.",
                  "Now: 908 cars displaced at the peak hour against 1,643 spaces "
                  "removed — measured in Yerevan.",
                  "Hourly licence-plate sweeps, 07:00–24:00, four working days."])
    return s


def v_supply_impact(prs):
    s = base_slide(prs, "Supply baseline and corridor impact",
                   "15,897 spaces surveyed; the conceptual design removes 85% of "
                   "on-corridor parking")
    xs, w = grid(2)
    tw = 2.28
    tiles = [("15,897", "Total spaces surveyed", INK),
             ("7,005", "On-street, on the corridors", INK),
             ("5,953", "Removed by the design", RED),
             ("1,052", "Retained or re-established", GREEN)]
    for i, (v, l, c) in enumerate(tiles):
        tile(s, L + (i % 2) * (tw + 0.18), ROW1 + (i // 2) * 1.26, tw, v, l,
             color=c, accent=c, h=1.16, vsize=18)
    add_img(s, "app_regulation.png", 5.10, ROW1, w=4.66)
    write(textbox(s, 5.10, 4.30, 4.66, 0.20),
          "Parking Regulation — yerevan-parking.vercel.app", size=7.5,
          color=GREY, italic=True)
    card(s, L, 4.08, 4.74, 0.72, "88% of it is free — and unorganised", None,
         accent=AMBER, bsize=9,
         bullets=["Only 20% has signage; 12% has markings."])
    footnote(s, "Removal is inherent to the central-median cross-section — "
                "1,836 / 3,310 / 807 across Corridors 1–3.", y=4.88)
    return s


def v_demand(prs):
    s = base_slide(prs, "A saturated kerb — but a short-stay one",
                   "69% of zones sit above the 85% efficiency threshold at the "
                   "peak, yet almost two-thirds of vehicles stay an hour or less")
    write(textbox(s, L, ROW1 - 0.18, 4.8, 0.20), "Peak occupancy by surveyed area",
          size=9, bold=True, color=INK)
    add_img(s, "chart_occupancy.png", L, ROW1 + 0.04, w=4.86)
    write(textbox(s, 5.35, ROW1 - 0.18, 3.45, 0.20), "How long vehicles stay",
          size=9, bold=True, color=INK)
    add_img(s, "chart_stay.png", 5.35, ROW1 + 0.10, w=4.41)
    card(s, 5.35, 2.72, 4.41, 1.95, "What this rules out", None, accent=MED_T,
         hcolor=MED_T, bsize=9, bspace=6,
         bullets=["Visitor caps address a problem the data does not show — "
                  "turnover is already 5–11 vehicles per space per day.",
                  "All-day parking is only ~6%; overnight residential ~7%.",
                  "Pricing, not maximum-stay enforcement, is the binding lever."])
    footnote(s, "Occupancy above 100% is chronic overflow — more vehicles present "
                "than legal striped capacity. Only ~0.2% parked illegally: the "
                "issue is saturation, not illegality.", y=4.82)
    return s


def v_displacement(prs):
    s = base_slide(prs, "Displacement: measured, not inferred",
                   "Peak displaced demand is only 55% of the spaces removed — "
                   "but re-absorption depends on opening gated yards")
    add_img(s, "chart_displacement.png", L, ROW1 + 0.06, w=5.05)
    write(textbox(s, L, 3.06, 5.05, 0.20),
          "Where those 908 cars can go", size=9, bold=True, color=INK)
    add_img(s, "chart_absorption.png", L, 3.28, w=5.05)
    tile(s, 5.55, ROW1, 2.00, "55%", "Displaced vs. removed", color=RED,
         accent=RED, h=1.30, vsize=24)
    tile(s, 7.76, ROW1, 2.00, "92%", "Of absorptive capacity is behind gates",
         color=HIGH_T, accent=HIGH_T, h=1.30, vsize=24)
    card(s, 5.55, 3.06, 4.21, 1.60, "The condition attached", None, accent=HIGH_T,
         hcolor=HIGH_T, bsize=9, bspace=6,
         bullets=["Nearby streets absorb just 43% — 14% in Kentron, 0% in Komitas.",
                  "The rest needs off-street capacity that is still gated, "
                  "counted as if already open.",
                  "So: open and manage the yards BEFORE removing kerb."])
    footnote(s, "Surveyed areas only — a subset of the corridor-wide 5,953 removals.",
             y=4.82)
    return s


# --- superseded text-heavy tight slides (kept for reference) --------------------
def t_supply_impact(prs):
    s = base_slide(prs, "Supply baseline and corridor impact",
                   "15,897 spaces surveyed; the conceptual design removes 85% of "
                   "on-corridor parking")
    xs, w = grid(4)
    tiles = [("15,897", "Total spaces surveyed (34.6 km + 100 m buffer)", INK),
             ("7,005", "On-street, on the corridors themselves", INK),
             ("5,953", "On-street spaces removed by the design", RED),
             ("1,052", "Spaces retained or re-established", GREEN)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    stacked_zone_bar(s, 2.92)
    card(s, L, 4.02, CW, 0.72,
         "88% of corridor parking is free — and barely organised", None,
         accent=AMBER,
         bullets=["Only 20% has signage and 12% has markings. Removal is inherent "
                  "to the central-median cross-section, and is consistent with the "
                  "Mott MacDonald (2022) strategy and the Yerevan Sustainable Urban "
                  "Transport Strategy."])
    footnote(s, "Totals corrected June 2026 (off-street 7,479 → 7,035; total "
                "16,341 → 15,897). Removal splits 1,836 / 3,310 / 807 across "
                "Corridors 1–3.")
    return s


def t_displacement(prs):
    s = base_slide(prs, "Displacement: measured, not inferred",
                   "Peak displaced demand is only 55% of the spaces removed — but "
                   "re-absorption depends on opening gated yards")
    xs, w = grid(4)
    tiles = [("1,643", "Spaces removed in the surveyed areas", GREY),
             ("908", "Cars displaced at the peak hour", RED),
             ("55%", "Displaced demand vs. supply removed", GREEN),
             ("92%", "Of absorptive capacity sits behind gates", HIGH_T)]
    for x, (v, l, c) in zip(xs, tiles):
        tile(s, x, ROW1, w, v, l, color=c, accent=c)
    xs2, w2 = grid(2)
    card(s, xs2[0], 3.05, w2, 1.55, "Measured in Yerevan", None, accent=GREEN,
         bsize=9, bspace=6,
         bullets=["The European BHLS finding — displaced demand falls well below "
                  "supply removed — is now confirmed by local measurement, not "
                  "borrowed from precedent.",
                  "Not every removed space needs replacing one-for-one: users shift "
                  "mode, retime, or were parked informally."])
    card(s, xs2[1], 3.05, w2, 1.55, "The condition attached", None, accent=HIGH_T,
         hcolor=HIGH_T, bsize=9, bspace=6,
         bullets=["Surviving on-street capacity absorbs only 43% — 14% in "
                  "Kentron, 0% in Komitas, up to 97% in Malatia-Sebastia.",
                  "The balance is closed by off-street, but 92% of that is gated "
                  "yards counted as if already open.",
                  "Open and manage the yards BEFORE removing kerb."])
    footnote(s, "Surveyed areas only — a subset of the corridor-wide 5,953 "
                "removals. Three of the six surveyed off-street facilities are "
                "already at or over capacity at the peak hour.")
    return s


# --- build --------------------------------------------------------------------
def build(path, slides):
    prs = new_deck()
    for fn in slides:
        fn(prs)
    prs.save(path)
    print("  %-46s %2d slides" % (path.split("/")[-1], len(slides)))


if __name__ == "__main__":
    print("Building parking slides for the final ADB mission presentation:")
    ensure_assets()
    build(OUT_FULL, [s_whats_new, s_approach, s_supply, s_occupancy, s_duration,
                     s_impact, s_displacement, s_yards, s_verdicts, s_zones,
                     s_conclusions])
    build(OUT_TIGHT, [v_measured, v_supply_impact, v_demand, v_displacement,
                      s_zones])
    print("Done.")
